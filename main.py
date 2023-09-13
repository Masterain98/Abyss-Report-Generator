from pptx import Presentation
from datetime import datetime
import homa
import pptxutils
import pytz
import os
import win32com.client
from pathlib import Path


def get_replacement_dict():
    return_dict = {}
    new_return_dict = {}

    # ScheduleId
    overview = homa.get_overview()
    return_dict["s"] = overview["ScheduleId"]
    return_dict["record-total"] = overview["RecordTotal"]
    return_dict["abyss-total"] = overview["SpiralAbyssTotal"]
    return_dict["abyss-passed"] = overview["SpiralAbyssPassed"]
    return_dict["pass-rate"] = f'{(overview["SpiralAbyssPassed"] / overview["SpiralAbyssTotal"] * 100):.2f}%'
    return_dict["full-star"] = overview["SpiralAbyssFullStar"]
    return_dict["fs-rate"] = f'{(overview["SpiralAbyssFullStar"] / overview["SpiralAbyssTotal"] * 100):.2f}%'

    # Title
    current_time = datetime.now(tz=pytz.timezone("Asia/Shanghai"))
    return_dict["y"] = current_time.year
    return_dict["m"] = current_time.month
    return_dict["date"] = current_time.strftime("%Y/%m/%d") + " UTC+8"
    if current_time.day < 15:
        return_dict["h"] = "上"
    else:
        return_dict["h"] = "下"

    # 使用率和出场率
    floor_12_top_20_utilization_rate = homa.get_floor_12_top_20_utilization_rate()
    floor_12_top_20_attendance_rate = homa.get_floor_12_top_20_attendance_rate()
    for i in range(20):
        if i + 1 < 10:
            i_text = "0" + str(i + 1)
        else:
            i_text = str(i + 1)
        return_dict[f"sy-{i_text}-png"] = list(floor_12_top_20_utilization_rate.keys())[i]
        return_dict[f"sy-{i_text}"] = list(floor_12_top_20_utilization_rate.values())[i]
        return_dict[f"cc-{i_text}-png"] = list(floor_12_top_20_attendance_rate.keys())[i]
        return_dict[f"cc-{i_text}"] = list(floor_12_top_20_attendance_rate.values())[i]

    # 队伍组合
    floor_12_top_3_team_combination = homa.get_floor_12_top_3_team_combination()
    for i in range(3):
        return_dict[f"u{str(i + 1)}-1-png"] = floor_12_top_3_team_combination["up"][i][f"u{str(i + 1)}-1-png"]
        return_dict[f"u{str(i + 1)}-2-png"] = floor_12_top_3_team_combination["up"][i][f"u{str(i + 1)}-2-png"]
        return_dict[f"u{str(i + 1)}-3-png"] = floor_12_top_3_team_combination["up"][i][f"u{str(i + 1)}-3-png"]
        return_dict[f"u{str(i + 1)}-4-png"] = floor_12_top_3_team_combination["up"][i][f"u{str(i + 1)}-4-png"]
        return_dict[f"u{str(i + 1)}"] = floor_12_top_3_team_combination["up"][i]["count"]
        return_dict[f"d{str(i + 1)}-1-png"] = floor_12_top_3_team_combination["down"][i][f"d{str(i + 1)}-1-png"]
        return_dict[f"d{str(i + 1)}-2-png"] = floor_12_top_3_team_combination["down"][i][f"d{str(i + 1)}-2-png"]
        return_dict[f"d{str(i + 1)}-3-png"] = floor_12_top_3_team_combination["down"][i][f"d{str(i + 1)}-3-png"]
        return_dict[f"d{str(i + 1)}-4-png"] = floor_12_top_3_team_combination["down"][i][f"d{str(i + 1)}-4-png"]
        return_dict[f"d{str(i + 1)}"] = floor_12_top_3_team_combination["down"][i]["count"]

    for k, v in return_dict.items():
        new_return_dict["{{" + k + "}}"] = v

    return new_return_dict


def replace_text(self, replacements: dict, shapes: list):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, str(replacement))
                            cell.text = new_text


def ppt2png(ppt_path, output_path):
    ppt_app = win32com.client.Dispatch('PowerPoint.Application')
    ppt = ppt_app.Presentations.Open(ppt_path)
    page_number = 1
    for slide in ppt.Slides:
        slide.Export(output_path + f'/Page {page_number}.png', 'png', 3840, 2160)
        page_number += 1
    ppt.SaveAs(output_path + '/abyss-report.pdf', 32)
    ppt_app.Quit()


def main():
    os.makedirs("output", exist_ok=True)
    r_dict = get_replacement_dict()
    prs = Presentation('abyss-template.pptx')

    slides = [slide for slide in prs.slides]
    for slide in slides:
        replace_text(slide, r_dict, slide.shapes)

        for shape in slide.shapes:
            alt_text = pptxutils.shape_alt_text(shape)
            if alt_text in r_dict.keys():
                new_shape = slide.shapes.add_picture(
                    r_dict[alt_text], shape.left, shape.top, shape.width, shape.height
                )
                old_pic = shape.element
                new_pic = new_shape.element
                old_pic.addnext(new_pic)
                old_pic.getparent().remove(old_pic)
    prs.save('abyss-report.pptx')
    ppt2png(str(Path().absolute())+"/abyss-report.pptx", str(Path().absolute())+"/output/")


if __name__ == "__main__":
    main()
